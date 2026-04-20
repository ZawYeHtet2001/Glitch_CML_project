"""
Final crit deck for Interactive Memory Machine (v2.5).
Tuesday 2026-04-21 evening crit, SUTD CML.

Run: python generate_crit_deck.py
Output: IMM_Final_Crit.pptx (~14 slides)

Deliberately lean on real artifacts (public/archetypes + Test_imgs) rather
than text bullets. Each slide is meant to support a spoken point — not
replace it.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ========== Colors (match app v2.5 palette) ==========
BG_DEEP   = RGBColor(0x07, 0x05, 0x02)   # near-black with warm undertone
BG_PANEL  = RGBColor(0x12, 0x0a, 0x02)   # CRT amber-black
CARD      = RGBColor(0x1a, 0x13, 0x06)
AMBER     = RGBColor(0xff, 0xb3, 0x47)   # phosphor hot
AMBER_DIM = RGBColor(0xc8, 0x8a, 0x34)
INK_DARK  = RGBColor(0x2b, 0x27, 0x21)
BONE      = RGBColor(0xd7, 0xcc, 0xb4)
FG        = RGBColor(0xed, 0xed, 0xed)
MUTED     = RGBColor(0x8a, 0x82, 0x6a)
DIM       = RGBColor(0x55, 0x4c, 0x3a)
RED       = RGBColor(0xff, 0x5a, 0x5a)
GREEN     = RGBColor(0x6c, 0xcf, 0x7f)
WHITE     = RGBColor(0xff, 0xff, 0xff)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ARCHETYPE_DIR = "public/archetypes"
OUTPUT_DIR = "Test_imgs"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ========== Helpers ==========
def set_bg(slide, color=BG_DEEP):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tx(slide, left, top, width, height, text, *, size=14, color=FG,
       bold=False, align=PP_ALIGN.LEFT, font="Courier New", spacing=1.1,
       italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0);  tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(spacing)
    return tf


def bullets(slide, left, top, width, height, items, *, size=13, color=FG,
            font="Courier New"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(5)
    return tf


def rect(slide, left, top, width, height, *, fill=CARD, border=INK_DARK,
         border_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = border
        s.line.width = Pt(border_w)
    s.shadow.inherit = False
    return s


def accent_line(slide, left, top, width, color=AMBER):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def chassis_frame(slide, n, total):
    """Thin amber top + bottom strips, matches the app's lean chrome."""
    # Top strip
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.34))
    top.fill.solid(); top.fill.fore_color.rgb = BG_PANEL
    top.line.color.rgb = AMBER_DIM
    top.line.width = Pt(0.75)
    top.shadow.inherit = False
    tx(slide, Inches(0.25), Inches(0.05), Inches(7), Inches(0.25),
       "INTERACTIVE MEMORY MACHINE  ·  IMM-0497  ·  v2.5  ·  CH.01",
       size=9, color=AMBER_DIM)
    tx(slide, Inches(7.5), Inches(0.05), Inches(5.6), Inches(0.25),
       f"SLIDE {n:02d} / {total:02d}  ·  FINAL CRIT  ·  2026-04-21",
       size=9, color=AMBER_DIM, align=PP_ALIGN.RIGHT)

    # Bottom strip
    bot = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.16), SLIDE_W, Inches(0.34))
    bot.fill.solid(); bot.fill.fore_color.rgb = BG_PANEL
    bot.line.color.rgb = AMBER_DIM
    bot.line.width = Pt(0.75)
    bot.shadow.inherit = False
    tx(slide, Inches(0.25), Inches(7.2), Inches(7), Inches(0.25),
       "TEAM 10  ·  SUTD  ·  CREATIVE MACHINE LEARNING",
       size=9, color=AMBER_DIM)
    tx(slide, Inches(7.5), Inches(7.2), Inches(5.6), Inches(0.25),
       "JW  /  NT  /  VP  /  ZH",
       size=9, color=AMBER_DIM, align=PP_ALIGN.RIGHT)


def slide_kicker(slide, stage, title):
    """Small amber kicker + big white title, beneath the chassis top strip."""
    tx(slide, Inches(0.5), Inches(0.55), Inches(5), Inches(0.35),
       stage, size=11, color=AMBER, bold=True)
    accent_line(slide, Inches(0.5), Inches(0.95), Inches(1.8))
    tx(slide, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.7),
       title, size=26, color=FG, bold=True, font="Georgia")


def add_image(slide, path, left, top, width, height):
    """Safe image add — returns None if the file is unreadable."""
    if not os.path.exists(path):
        return None
    try:
        return slide.shapes.add_picture(path, left, top, width=width, height=height)
    except Exception:
        return None


TOTAL = 14


# ============================================================
# SLIDE 1 — Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 1, TOTAL)

# Large amber "beam" sweep effect — decorative
rect(s, Inches(0.5), Inches(1.15), Inches(2.5), Pt(3), fill=AMBER, border=None)

tx(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
   "◆ TRANSMISSION :: SUBCONSCIOUS CHANNEL OPEN",
   size=12, color=AMBER, bold=True)

tx(s, Inches(0.5), Inches(2.4), Inches(12.3), Inches(1.6),
   "INTERACTIVE",
   size=72, color=FG, bold=True, font="Georgia")
tx(s, Inches(0.5), Inches(3.5), Inches(12.3), Inches(1.6),
   "MEMORY MACHINE",
   size=72, color=AMBER, bold=True, font="Georgia")

accent_line(s, Inches(0.5), Inches(5.1), Inches(3.5))

tx(s, Inches(0.5), Inches(5.25), Inches(12), Inches(0.5),
   "Glitch as the meaningful distortion of architectural memory.",
   size=18, color=FG, italic=True, font="Georgia")

tx(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.4),
   "Creative Machine Learning  //  Term 6  //  SUTD",
   size=13, color=AMBER_DIM)
tx(s, Inches(0.5), Inches(6.55), Inches(12), Inches(0.4),
   "Team 10 — Zaw Ye Htet  ·  Jiawen  ·  Naing Thway  ·  Vrinda",
   size=12, color=MUTED)


# ============================================================
# SLIDE 2 — The Question
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 2, TOTAL)
slide_kicker(s, "01 // QUESTION",
             "What if architecture could be recalled the way we actually remember it?")

tx(s, Inches(0.5), Inches(2.6), Inches(8.0), Inches(1.2),
   '"Space is not remembered as a photograph.\n'
   'It tilts. It omits. It compresses.\n'
   'It fuses with other spaces that no longer exist."',
   size=20, color=AMBER, italic=True, font="Georgia", spacing=8)

bullets(s, Inches(0.5), Inches(5.0), Inches(12), Inches(1.5),
        [
            "// Memory is an unreliable archivist — and that unreliability is architectural.",
            "// The subconscious acts as an editor: warping scale, dissolving edges, merging bodies.",
            "// We wanted to externalise that process — make it visible, interactive, and 3D-printable.",
        ], size=13, color=FG)


# ============================================================
# SLIDE 3 — Concept: Glitch Reframed
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 3, TOTAL)
slide_kicker(s, "02 // POSITION", "Glitch is not aesthetic.")

tx(s, Inches(0.5), Inches(2.2), Inches(7), Inches(0.5),
   "It is how the subconscious edits architecture.",
   size=22, color=AMBER, italic=True, font="Georgia")

bullets(s, Inches(0.5), Inches(3.3), Inches(6.8), Inches(3.5),
        [
            "- We reject the 'pixel-artifact' reading of glitch.",
            "",
            "- In this project, glitch is meaningful distortion:",
            "    · Clarity dissolves, edges blur into neighbours.",
            "    · Completeness fails, volumes carve into voids.",
            "    · Stability tilts, mass suspends off-axis.",
            "    · Misassociation fuses rooms that never met.",
            "",
            "- The output must remain physically plausible — a sculpture,",
            "  not an abstract composition.",
        ], size=13, color=FG)

# Right: before / after diagram
right_x = Inches(7.8)
rect(s, right_x, Inches(2.2), Inches(5.05), Inches(4.5),
     fill=CARD, border=INK_DARK)
tx(s, right_x + Inches(0.3), Inches(2.35), Inches(4.5), Inches(0.3),
   "ORDERED  ->  RECALLED", size=10, color=AMBER, bold=True)

# clean cube (abstract)
cx = right_x + Inches(0.8); cy = Inches(3.2)
rect(s, cx, cy, Inches(1.4), Inches(1.4), fill=BG_PANEL, border=MUTED)
rect(s, cx + Inches(0.2), cy - Inches(0.2), Inches(1.4), Inches(0.2),
     fill=AMBER_DIM, border=MUTED)
tx(s, cx - Inches(0.1), cy + Inches(1.55), Inches(1.6), Inches(0.3),
   "AS BUILT", size=10, color=MUTED, align=PP_ALIGN.CENTER)

# arrow
tx(s, right_x + Inches(2.35), cy + Inches(0.55), Inches(0.5), Inches(0.5),
   "->", size=28, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# distorted cube — scattered shards
dx = right_x + Inches(3.05); dy = Inches(3.0)
rect(s, dx, dy + Inches(0.1), Inches(0.9), Inches(0.95),
     fill=BG_PANEL, border=AMBER_DIM)
rect(s, dx + Inches(1.0), dy + Inches(0.4), Inches(0.55), Inches(0.65),
     fill=BG_PANEL, border=AMBER_DIM)
rect(s, dx + Inches(0.3), dy + Inches(1.2), Inches(1.2), Inches(0.45),
     fill=BG_PANEL, border=AMBER)
rect(s, dx + Inches(1.5), dy + Inches(1.25), Inches(0.35), Inches(0.35),
     fill=AMBER_DIM, border=AMBER_DIM)
rect(s, dx - Inches(0.15), dy + Inches(0.3), Inches(0.22), Inches(0.22),
     fill=AMBER, border=AMBER)
tx(s, dx - Inches(0.2), cy + Inches(1.55), Inches(2), Inches(0.3),
   "AS RECALLED", size=10, color=AMBER, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4 — The Machine (one-liner + screenshot proxy)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 4, TOTAL)
slide_kicker(s, "03 // SYSTEM", "The Interactive Memory Machine.")

tx(s, Inches(0.5), Inches(2.1), Inches(12), Inches(0.9),
   "A browser-based laboratory where a subject maps their own recall to geometric distortion.",
   size=16, color=AMBER_DIM, italic=True, font="Georgia", spacing=8)

# Four quadrants summarising the machine
blocks = [
    ("INPUT",       "Subconscious spatial recall\n(free text, <=600 chars)"),
    ("ANALYSIS",    "Claude extracts keywords, scores\n12 distortion parameters, picks archetype"),
    ("MAPPING",     "Subject wires their keywords\nto operation nodes on a canvas"),
    ("ARTIFACT",    "fal.ai generates 2D image,\nthen lifts it into a 3D printable mesh"),
]

bx = Inches(0.5); by = Inches(3.3)
bw = Inches(3.0); bh = Inches(2.9); bgap = Inches(0.15)
for i, (h, body) in enumerate(blocks):
    x = bx + i * (bw + bgap)
    rect(s, x, by, bw, bh, fill=BG_PANEL, border=AMBER_DIM)
    tx(s, x + Inches(0.25), by + Inches(0.25), bw - Inches(0.5), Inches(0.45),
       f"0{i+1}  {h}", size=13, color=AMBER, bold=True)
    accent_line(s, x + Inches(0.25), by + Inches(0.78),
                bw - Inches(0.5), color=AMBER_DIM)
    tx(s, x + Inches(0.25), by + Inches(1.0), bw - Inches(0.5), Inches(1.8),
       body, size=12, color=FG, spacing=6)

tx(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.4),
   "// No existing catalogue. The subject is the author. The glitch is the testimony.",
   size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5 — Pipeline (10 stages)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 5, TOTAL)
slide_kicker(s, "04 // PIPELINE", "Ten stages from recalled text to exported mesh.")

pipeline = [
    ("01", "INPUT",       "recall text"),
    ("02", "ANALYZE",     "Claude — keywords,\n12 scores, archetype"),
    ("03", "CANVAS",      "wire keywords\nto operations"),
    ("04", "DIRECT",      "Claude rewrites\nmechanical prompt"),
    ("05", "GENERATE",    "fal.ai — Recraft V4\n/ Nano / Ideogram"),
    ("06", "NAME",        "Claude returns\n{name, explanation}"),
    ("07", "VIEW 2D",     "crisp image + title\n+ interpretation"),
    ("08", "RECONSTRUCT", "fal.ai — Trellis\nor Hunyuan3D"),
    ("09", "ORBIT",       "R3F viewer\nSTUDIO / MOODY"),
    ("10", "EXPORT",      "GLB + OBJ + STL\nclient-side"),
]

# Two rows of 5
row_y = [Inches(2.4), Inches(4.75)]
cw = Inches(2.3); ch = Inches(2.0); cgap = Inches(0.2)
for i, (num, title, desc) in enumerate(pipeline):
    row = i // 5
    col = i % 5
    x = Inches(0.5) + col * (cw + cgap)
    y = row_y[row]

    # Alternate fill — AI stages slightly amber, user stages plain
    ai_stage = num in ("02", "04", "05", "06", "08")
    border_c = AMBER if ai_stage else AMBER_DIM
    rect(s, x, y, cw, ch, fill=BG_PANEL, border=border_c)
    tx(s, x + Inches(0.2), y + Inches(0.15), Inches(1), Inches(0.4),
       num, size=20, color=AMBER if ai_stage else FG, bold=True)
    tx(s, x + Inches(0.2), y + Inches(0.7), cw - Inches(0.4), Inches(0.35),
       title, size=12, color=FG, bold=True)
    tx(s, x + Inches(0.2), y + Inches(1.1), cw - Inches(0.4), Inches(1),
       desc, size=9, color=AMBER_DIM, spacing=1.0)

    if col < 4:
        tx(s, x + cw, y + Inches(0.75), cgap, Inches(0.4),
           ">", size=16, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

tx(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
   "AMBER = MODEL INFERENCE (CLAUDE / FAL.AI)   ·   WHITE = USER ACTION",
   size=9, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6 — Archetypes (7 tonal starters with real images)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 6, TOTAL)
slide_kicker(s, "05 // ARCHETYPES",
             "Seven tonal starters — Claude classifies each recall into one.")

archetypes = [
    ("ORGANIC",     "organic.jpg",     "Calm / nostalgic",   "Smooth biomorphic curves"),
    ("CRYSTALLINE", "crystalline.jpg", "Tense / anxious",    "Angular faceted planes"),
    ("TWISTED",     "twisted.jpg",     "Disorienting",       "Coiling spirals"),
    ("SKELETAL",    "skeletal.jpg",    "Vulnerable",         "Thin branching"),
    ("MONOLITHIC",  "monolithic.jpg",  "Oppressive",         "Dense compact mass"),
    ("FRAGMENTED",  "fragmented.jpg",  "Chaotic / traumatic","Scattered shards"),
    ("NESTED",      "nested.jpg",      "Introspective",      "Concentric shells"),
]

card_w = Inches(1.70); card_h = Inches(4.25)
gap = Inches(0.12); x0 = Inches(0.5); y0 = Inches(2.3)

for i, (name, img, mood, form) in enumerate(archetypes):
    x = x0 + i * (card_w + gap)
    rect(s, x, y0, card_w, card_h, fill=BG_PANEL, border=AMBER_DIM)
    img_path = os.path.join(ARCHETYPE_DIR, img)
    if os.path.exists(img_path):
        try:
            s.shapes.add_picture(img_path,
                                 x + Inches(0.1), y0 + Inches(0.1),
                                 width=card_w - Inches(0.2),
                                 height=Inches(1.85))
        except Exception:
            pass

    tx(s, x + Inches(0.1), y0 + Inches(2.08), card_w - Inches(0.2), Inches(0.35),
       name, size=11, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(0.1), y0 + Inches(2.5), card_w - Inches(0.2), Inches(0.5),
       mood, size=9, color=FG, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(0.1), y0 + Inches(3.1), card_w - Inches(0.2), Inches(1),
       form, size=9, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.0)

tx(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.3),
   "// The archetype seeds the shape-language. The 12 operation nodes deform it.",
   size=11, color=AMBER_DIM, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — 12 Operation Nodes
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 7, TOTAL)
slide_kicker(s, "06 // OPERATIONS",
             "Twelve nodes — three families of distortion.")

groups = [
    ("SPATIAL (0..1)", [
        ("CLARITY",       "Boundary diffusion"),
        ("COMPLETENESS",  "Mass subtraction"),
        ("STABILITY",     "Gravitational shift"),
        ("MISASSOC.",     "Cross-morph fusion"),
        ("VULNERABILITY", "Interior reveal"),
        ("INTIMACY",      "Scale compression"),
    ]),
    ("EXPERIENCE (bipolar)", [
        ("TEMPERATURE", "Freeze <-> Melt"),
        ("PRESSURE",    "Smooth <-> Shatter"),
        ("LUMINOSITY",  "Shadow <-> Radiance"),
    ]),
    ("APPEARANCE (0..1)", [
        ("MATERIAL", "Substance identity"),
        ("TEXTURE",  "Surface grain"),
        ("COLOR",    "Chromatic palette"),
    ]),
]

col_x = [Inches(0.5), Inches(5.0), Inches(9.5)]
col_w = [Inches(4.3), Inches(4.3), Inches(3.35)]
y_top = Inches(2.3)

for ci, (gname, nodes) in enumerate(groups):
    x = col_x[ci]; w = col_w[ci]
    # group header
    rect(s, x, y_top, w, Inches(0.55), fill=AMBER_DIM, border=AMBER)
    tx(s, x + Inches(0.2), y_top + Inches(0.13), w - Inches(0.4), Inches(0.35),
       gname, size=12, color=INK_DARK, bold=True)

    for i, (name, desc) in enumerate(nodes):
        ny = y_top + Inches(0.7) + i * Inches(0.65)
        rect(s, x, ny, w, Inches(0.58), fill=BG_PANEL, border=AMBER_DIM)
        tx(s, x + Inches(0.18), ny + Inches(0.08), w - Inches(0.4), Inches(0.3),
           name, size=11, color=AMBER, bold=True)
        tx(s, x + Inches(0.18), ny + Inches(0.32), w - Inches(0.4), Inches(0.25),
           desc, size=10, color=FG)

tx(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3),
   "// Connection COUNT and keyword CONTENT together modulate the intensity of each operation.",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 8 — Node Canvas: the testimony
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 8, TOTAL)
slide_kicker(s, "07 // AUTHORSHIP",
             "The node canvas — the subject maps their own distortion.")

# Canvas mock
canvas_x = Inches(0.5); canvas_y = Inches(2.2)
canvas_w = Inches(8.2); canvas_h = Inches(4.3)
rect(s, canvas_x, canvas_y, canvas_w, canvas_h, fill=BG_PANEL, border=AMBER_DIM)
tx(s, canvas_x + Inches(0.2), canvas_y + Inches(0.1),
   Inches(3), Inches(0.3),
   "KEYWORDS", size=9, color=AMBER_DIM, bold=True)
tx(s, canvas_x + Inches(5.4), canvas_y + Inches(0.1),
   Inches(3), Inches(0.3),
   "OPERATIONS", size=9, color=AMBER_DIM, bold=True)

kws = [
    ("OBJ", "corridor"), ("OBJ", "stove"), ("PER", "mother"),
    ("SPA", "narrow"), ("EXP", "cold"), ("EVT", "shifting"),
    ("OBJ", "kettle"), ("SPA", "high"),
]
kw_x = canvas_x + Inches(0.25)
kw_y0 = canvas_y + Inches(0.55)
cat_col = {"OBJ": FG, "PER": AMBER, "EVT": GREEN, "SPA": AMBER_DIM, "EXP": RED}
for i, (cat, word) in enumerate(kws):
    y = kw_y0 + i * Inches(0.42)
    rect(s, kw_x, y, Inches(2.1), Inches(0.36), fill=CARD, border=AMBER_DIM)
    tx(s, kw_x + Inches(0.1), y + Inches(0.08), Inches(0.4), Inches(0.2),
       cat, size=8, color=cat_col.get(cat, FG), bold=True)
    tx(s, kw_x + Inches(0.55), y + Inches(0.08), Inches(1.4), Inches(0.2),
       word, size=10, color=FG)
    port = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              kw_x + Inches(2.02), y + Inches(0.12),
                              Inches(0.14), Inches(0.14))
    port.fill.solid(); port.fill.fore_color.rgb = AMBER
    port.line.color.rgb = AMBER; port.line.width = Pt(0.5)

ops = [
    ("CLARITY", 0.72), ("COMPLETENESS", 0.35), ("STABILITY", 0.81),
    ("MISASSOC.", 0.28), ("VULNERABILITY", 0.55), ("INTIMACY", 0.67),
]
op_x = canvas_x + Inches(5.4); op_w = Inches(2.55)
op_y0 = canvas_y + Inches(0.55)
for i, (name, score) in enumerate(ops):
    y = op_y0 + i * Inches(0.57)
    rect(s, op_x, y, op_w, Inches(0.5), fill=CARD, border=AMBER_DIM)
    port = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              op_x - Inches(0.08), y + Inches(0.18),
                              Inches(0.16), Inches(0.16))
    port.fill.solid(); port.fill.fore_color.rgb = AMBER
    port.line.color.rgb = AMBER; port.line.width = Pt(0.5)
    tx(s, op_x + Inches(0.13), y + Inches(0.03),
       op_w - Inches(0.7), Inches(0.25),
       name, size=10, color=AMBER, bold=True)
    tx(s, op_x + op_w - Inches(0.6), y + Inches(0.03),
       Inches(0.55), Inches(0.25),
       f"{score:.2f}", size=10, color=FG, align=PP_ALIGN.RIGHT)
    bar_bg_w = Inches(2.2)
    rect(s, op_x + Inches(0.13), y + Inches(0.33), bar_bg_w, Inches(0.06),
         fill=INK_DARK, border=INK_DARK)
    rect(s, op_x + Inches(0.13), y + Inches(0.33),
         Inches(2.2 * score), Inches(0.06),
         fill=AMBER, border=AMBER)

# wires
connections = [(0, 0), (1, 5), (3, 2), (4, 5), (5, 3), (6, 0), (7, 2)]
for (ki, oi) in connections:
    sx = kw_x + Inches(2.09)
    sy = kw_y0 + ki * Inches(0.42) + Inches(0.19)
    ex = op_x - Inches(0.01)
    ey = op_y0 + oi * Inches(0.57) + Inches(0.26)
    c = s.shapes.add_connector(MSO_CONNECTOR.CURVE, sx, sy, ex, ey)
    c.line.color.rgb = AMBER
    c.line.width = Pt(1.1)

tx(s, canvas_x + Inches(3.0), canvas_y + Inches(2.0), Inches(2.4), Inches(0.5),
   "DRAG  ->  WIRE  ->  DEFORM",
   size=12, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# Right copy
tx(s, Inches(9.0), Inches(2.3), Inches(4.0), Inches(0.5),
   "The wiring is the testimony.",
   size=18, color=AMBER, italic=True, bold=True, font="Georgia")
bullets(s, Inches(9.0), Inches(3.1), Inches(4.0), Inches(3.8),
        [
            "// Multiple keywords per node =",
            "  stronger pull on that operation.",
            "",
            "// Nothing is pre-decided. The same",
            "  recall produces different artifacts",
            "  based on wiring choices.",
            "",
            "// The canvas is the subject's",
            "  authored reading of their own memory.",
        ], size=12, color=FG)


# ============================================================
# SLIDE 9 — Art Director
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 9, TOTAL)
slide_kicker(s, "08 // TRANSLATION",
             "An art director closes the gap between spec and image.")

# Left: mechanical spec
mx = Inches(0.5); my = Inches(2.2); mw = Inches(4.0); mh = Inches(4.4)
rect(s, mx, my, mw, mh, fill=BG_PANEL, border=AMBER_DIM)
tx(s, mx + Inches(0.2), my + Inches(0.2), mw - Inches(0.4), Inches(0.3),
   "MECHANICAL PROMPT", size=10, color=AMBER, bold=True)
tx(s, mx + Inches(0.2), my + Inches(0.5), mw - Inches(0.4), Inches(0.3),
   "(numbers + tokens)", size=9, color=MUTED)
bullets(s, mx + Inches(0.2), my + Inches(1.0), mw - Inches(0.4), mh - Inches(1.2),
        [
            "archetype: CRYSTALLINE",
            "",
            "clarity 0.72",
            "completeness 0.35",
            "stability 0.81",
            "temperature 0.22",
            "pressure 0.66",
            "luminosity 0.41",
            "",
            "kw(corridor, narrow, high)",
            "  -> SPATIAL nodes",
            "kw(cold, shifting)",
            "  -> TEMPERATURE low",
            "",
            "+ 9 more parameters...",
        ], size=10, color=FG)

# Arrow
tx(s, Inches(4.55), Inches(4.1), Inches(0.5), Inches(0.5),
   "->", size=24, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# Centre: the director (Claude)
cx = Inches(5.2); cy = Inches(2.2)
cw = Inches(3.0); ch = Inches(4.4)
rect(s, cx, cy, cw, ch, fill=CARD, border=AMBER, border_w=1.5)
tx(s, cx + Inches(0.2), cy + Inches(0.2), cw - Inches(0.4), Inches(0.3),
   "ART DIRECTOR", size=10, color=AMBER, bold=True)
tx(s, cx + Inches(0.2), cy + Inches(0.5), cw - Inches(0.4), Inches(0.3),
   "Claude Haiku", size=11, color=FG, bold=True)

bullets(s, cx + Inches(0.2), cy + Inches(1.1), cw - Inches(0.4), ch - Inches(1.4),
        [
            "READS",
            "  · original recall",
            "  · archetype",
            "  · all 12 scores",
            "  · wiring pattern",
            "",
            "WRITES",
            "  · single cohesive",
            "    visual brief",
            "",
            "Physically plausible",
            "sculpture, clean",
            "background.",
        ], size=9, color=FG)

tx(s, Inches(8.25), Inches(4.1), Inches(0.5), Inches(0.5),
   "->", size=24, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# Right: visual brief
vx = Inches(8.85); vy = Inches(2.2)
vw = Inches(4.0); vh = Inches(4.4)
rect(s, vx, vy, vw, vh, fill=BG_PANEL, border=AMBER_DIM)
tx(s, vx + Inches(0.2), vy + Inches(0.2), vw - Inches(0.4), Inches(0.3),
   "VISUAL BRIEF", size=10, color=AMBER, bold=True)
tx(s, vx + Inches(0.2), vy + Inches(0.5), vw - Inches(0.4), Inches(0.3),
   "(sent to fal.ai)", size=9, color=MUTED)
tx(s, vx + Inches(0.2), vy + Inches(1.0), vw - Inches(0.4), vh - Inches(1.2),
   "A crystalline sculpture of\n"
   "a narrow corridor, faceted\n"
   "ice-like planes dissolving\n"
   "at the edges, carved with\n"
   "deep hollows where mass\n"
   "has been lost. Tilts gently\n"
   "off-axis. Cold palette,\n"
   "matte frost texture.\n"
   "Studio lighting, clean\n"
   "white background,\n"
   "physically plausible,\n"
   "single sculptural object.",
   size=11, color=FG, font="Georgia", spacing=8)

tx(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.3),
   "// Without the director, text-to-image models treat mechanical specs as noise.",
   size=11, color=AMBER_DIM, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10 — 2D -> 3D
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 10, TOTAL)
slide_kicker(s, "09 // DIMENSION",
             "From the 2D artifact to a holdable mesh.")

stages = [
    ("IMAGE",       "generated 2D\nartifact",        AMBER_DIM),
    ("RECONSTRUCT", "fal.ai\nTrellis / Hunyuan3D",    AMBER),
    ("MESH",        "GLB with\nPBR materials",        AMBER_DIM),
    ("VIEWER",      "React Three Fiber\nSTUDIO / MOODY lighting", AMBER),
]
sx0 = Inches(0.5); sy = Inches(2.3)
sw = Inches(2.9); sh = Inches(1.95); sgap = Inches(0.2)

for i, (title, desc, col) in enumerate(stages):
    x = sx0 + i * (sw + sgap)
    rect(s, x, sy, sw, sh, fill=BG_PANEL, border=col)
    tx(s, x + Inches(0.2), sy + Inches(0.25), sw - Inches(0.4), Inches(0.4),
       title, size=14, color=col, bold=True)
    tx(s, x + Inches(0.2), sy + Inches(0.85), sw - Inches(0.4), Inches(1),
       desc, size=11, color=FG, spacing=4)
    if i < 3:
        tx(s, x + sw, sy + Inches(0.65), sgap + Inches(0.05), Inches(0.5),
           ">", size=20, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# Model comparison
ty = Inches(4.65)
rect(s, Inches(0.5), ty, Inches(12.3), Inches(2.15), fill=BG_PANEL, border=AMBER_DIM)
tx(s, Inches(0.75), ty + Inches(0.15), Inches(11.5), Inches(0.3),
   "MODEL TOGGLE  —  selected per generation, user choice",
   size=10, color=AMBER, bold=True)

hdr_y = ty + Inches(0.55)
headers = [("MODEL", Inches(0.75)), ("fal.ai ID", Inches(2.9)),
           ("COST", Inches(6.4)), ("FIDELITY", Inches(7.8)),
           ("WHEN TO USE", Inches(9.7))]
for (h, hx) in headers:
    tx(s, hx, hdr_y, Inches(3), Inches(0.3),
       h, size=10, color=MUTED, bold=True)
rect(s, Inches(0.75), hdr_y + Inches(0.33),
     Inches(11.8), Pt(0.75), fill=AMBER_DIM, border=None)

r1y = hdr_y + Inches(0.5)
tx(s, Inches(0.75), r1y, Inches(2), Inches(0.3), "TRELLIS",
   size=11, color=AMBER, bold=True)
tx(s, Inches(2.9), r1y, Inches(3.5), Inches(0.3),
   "fal-ai/trellis", size=10, color=FG)
tx(s, Inches(6.4), r1y, Inches(1.4), Inches(0.3),
   "~$0.02", size=10, color=GREEN)
tx(s, Inches(7.8), r1y, Inches(1.9), Inches(0.3),
   "standard", size=10, color=FG)
tx(s, Inches(9.7), r1y, Inches(3), Inches(0.3),
   "fast iteration, default", size=10, color=MUTED)

r2y = r1y + Inches(0.45)
tx(s, Inches(0.75), r2y, Inches(2), Inches(0.3), "HUNYUAN3D V3",
   size=11, color=AMBER, bold=True)
tx(s, Inches(2.9), r2y, Inches(3.5), Inches(0.3),
   "fal-ai/hunyuan3d-v3/image-to-3d",
   size=10, color=FG)
tx(s, Inches(6.4), r2y, Inches(1.4), Inches(0.3),
   "~$0.16", size=10, color=RED)
tx(s, Inches(7.8), r2y, Inches(1.9), Inches(0.3),
   "PBR materials, 500k faces",
   size=10, color=FG)
tx(s, Inches(9.7), r2y, Inches(3), Inches(0.3),
   "crit / exhibition hero",
   size=10, color=MUTED)

tx(s, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.25),
   "// Exports: GLB (native) + OBJ + STL (three-stdlib, binary STL opens in Bambu / Cura).",
   size=9, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11 — Sample outputs (four real artifacts)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 11, TOTAL)
slide_kicker(s, "10 // ARTIFACTS",
             "Four spatial memories, processed into sculpture.")

# Pick 4 available PNGs. Order recent-first.
candidates = [
    "POLnuSNFtmhqKJgpBHNVR_image.png",
    "ewNn7JYUY0HplEGQtjJ6X_image.png",
    "cRGG4a3hK9LPpTFytY7yx_image.png",
    "-6qT3xffyunkkBc_LitGF_image.png",
]
labels = [
    ("FRAGMENTED",  "Subject · ZH  /  parking-garage recall"),
    ("MONOLITHIC",  "Subject · NT  /  factory-annex recall"),
    ("TWISTED",     "Subject · JW  /  stairway recall"),
    ("SKELETAL",    "Subject · VP  /  kitchen recall"),
]

gw = Inches(2.95); gh = Inches(3.8)
gap = Inches(0.15); gx0 = Inches(0.5); gy = Inches(2.3)

for i, (fname, (arch, caption)) in enumerate(zip(candidates, labels)):
    x = gx0 + i * (gw + gap)
    # frame
    rect(s, x, gy, gw, gh, fill=CARD, border=AMBER_DIM)
    path = os.path.join(OUTPUT_DIR, fname)
    if os.path.exists(path):
        try:
            s.shapes.add_picture(path, x + Inches(0.1), gy + Inches(0.1),
                                 width=gw - Inches(0.2),
                                 height=Inches(2.75))
        except Exception:
            pass
    tx(s, x + Inches(0.1), gy + Inches(2.95), gw - Inches(0.2), Inches(0.35),
       arch, size=11, color=AMBER, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(0.1), gy + Inches(3.35), gw - Inches(0.2), Inches(0.4),
       caption, size=9, color=MUTED, align=PP_ALIGN.CENTER)

tx(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.35),
   "// Each print began as a 600-char recall text. None are manually sculpted.",
   size=12, color=AMBER_DIM, italic=True, align=PP_ALIGN.CENTER, font="Georgia")
tx(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3),
   "Teammates 3D-printed selected artifacts for the exhibition.",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12 — Interface / retro-CRT aesthetic
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 12, TOTAL)
slide_kicker(s, "11 // INTERFACE",
             "The machine has to feel like a machine.")

bullets(s, Inches(0.5), Inches(2.3), Inches(6.5), Inches(4.3),
        [
            "// Retro-CRT terminal aesthetic.",
            "  Amber phosphor on black, VT323 + Share Tech Mono.",
            "  Scanline overlays on every active viewport.",
            "",
            "// Tactile controls.",
            "  Chunky labeled buttons with pulsing LEDs.",
            "  Terminal-style > and $ prompts for input.",
            "  Mechanical loader — elapsed counter, ASCII",
            "  progress bar, cycling cryptic serial log.",
            "",
            "// Idle decay (on input only).",
            "  Letter-fall, scan-line sweep, periodic TV",
            "  glitch — the machine 'dreams' when left alone.",
            "",
            "// Performance resilience.",
            "  Hardened fal / Anthropic clients (timeouts,",
            "  queue status streaming). WebGL context-loss",
            "  auto-recovery in the 3D viewport.",
        ], size=12, color=FG)

# Right: sample CRT panel readout (pure vector mock)
rx = Inches(7.5); ry = Inches(2.3)
rw = Inches(5.3); rh = Inches(4.55)
rect(s, rx, ry, rw, rh, fill=BG_PANEL, border=AMBER_DIM)
tx(s, rx + Inches(0.2), ry + Inches(0.15), Inches(3), Inches(0.3),
   "ANALYSIS / RESULT", size=9, color=AMBER_DIM, bold=True)

tx(s, rx + Inches(0.2), ry + Inches(0.55),
   rw - Inches(0.4), Inches(0.4),
   "> SUBJECT: ZH  ·  11 KEYWORDS  ·  BASE SHAPE: FRAGMENTED",
   size=12, color=AMBER, bold=True)

# Serial log mock
log_lines = [
    ">> RCV :: RAW SPATIAL STREAM 0x4A2F",
    ">> DECODE :: SYNAPTIC VECTORS",
    ">> SEGMENT :: SUBJECT TOKENS [8..15]",
    ">> CROSS-REF :: ARCHETYPE ENGINE",
    ">> RESOLVE :: ARCHETYPE CANDIDATE",
    ">> WRITE :: KEYWORD POOL  OK",
]
for i, line in enumerate(log_lines):
    opacity_col = AMBER if i == len(log_lines) - 1 else AMBER_DIM
    tx(s, rx + Inches(0.3), ry + Inches(1.2) + i * Inches(0.35),
       rw - Inches(0.5), Inches(0.3),
       line, size=11, color=opacity_col)

# ASCII progress
tx(s, rx + Inches(0.3), ry + Inches(3.45), rw - Inches(0.5), Inches(0.35),
   "[████████████████████████░░░░░░░░] 73.4%",
   size=12, color=AMBER, bold=True)

tx(s, rx + Inches(0.3), ry + Inches(3.9), rw - Inches(0.5), Inches(0.3),
   "T+00023.5S / ~48S        FRAME 0168",
   size=10, color=AMBER_DIM)


# ============================================================
# SLIDE 13 — Reflection / What It Proves
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 13, TOTAL)
slide_kicker(s, "12 // REFLECTION",
             "What this experiment proves — and what it leaves open.")

tx(s, Inches(0.5), Inches(2.2), Inches(12), Inches(0.5),
   "PROVEN",
   size=13, color=AMBER, bold=True)
bullets(s, Inches(0.5), Inches(2.7), Inches(12), Inches(2.0),
        [
            "- Subconscious spatial recall can be externalised as a structured,",
            "  authored distortion pattern without flattening its ambiguity.",
            "- ML can be a translator layer between clinical parameters and",
            "  generative models without obliterating subject authorship.",
            "- The glitch remains architectural — every artifact is physically",
            "  plausible and printable, not abstract 2D composition.",
        ], size=12, color=FG)

tx(s, Inches(0.5), Inches(4.8), Inches(12), Inches(0.5),
   "OPEN",
   size=13, color=AMBER, bold=True)
bullets(s, Inches(0.5), Inches(5.3), Inches(12), Inches(1.7),
        [
            "- Inter-subjective comparison — do memories of the same space",
            "  from different people converge or diverge under the machine?",
            "- Longitudinal recall — does the same subject's artifact drift",
            "  as their memory of the space mutates over time?",
            "- Architectural use — can this inform design briefs for spaces",
            "  meant to be forgotten, reassembled, or grieved?",
        ], size=12, color=FG)


# ============================================================
# SLIDE 14 — Close
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s, BG_DEEP)
chassis_frame(s, 14, TOTAL)

tx(s, Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.2),
   "THANK YOU", size=72, color=FG, bold=True, font="Georgia")
accent_line(s, Inches(0.5), Inches(4.0), Inches(3))
tx(s, Inches(0.5), Inches(4.15), Inches(12.3), Inches(0.6),
   "Questions, critique, misremembered corridors welcomed.",
   size=16, color=AMBER, italic=True, font="Georgia")

tx(s, Inches(0.5), Inches(5.4), Inches(12), Inches(0.4),
   "TEAM 10  //  Creative Machine Learning  //  Term 6  //  SUTD",
   size=12, color=AMBER_DIM)
tx(s, Inches(0.5), Inches(5.8), Inches(12), Inches(0.4),
   "Jiawen  ·  Naing Thway  ·  Vrinda  ·  Zaw Ye Htet (primary dev)",
   size=11, color=MUTED)
tx(s, Inches(0.5), Inches(6.3), Inches(12), Inches(0.4),
   "Repo  github.com/ZawYeHtet2001/Glitch_CML_project",
   size=10, color=MUTED)
tx(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
   "◆ END TRANSMISSION",
   size=10, color=AMBER_DIM, bold=True)


# ============================================================
# Save
# ============================================================
out = "IMM_Final_Crit.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
