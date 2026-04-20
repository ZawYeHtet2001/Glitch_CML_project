"""
Generate v2.4 team update PowerPoint for Interactive Memory Machine.
Covers full 2D -> 3D pipeline, 12 nodes, 7 archetypes, art director,
mechanical/subconscious UI theme, and final crit plan.

Run: python generate_slides_v2_4.py
Output: IMM_v2_4_Team_Update.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# --- Colors (match app theme) ---
BG = RGBColor(0x0A, 0x0A, 0x0A)
CARD = RGBColor(0x14, 0x14, 0x14)
CARD_HI = RGBColor(0x1C, 0x1C, 0x1C)
ACCENT = RGBColor(0xD4, 0xA0, 0x17)
ACCENT_DIM = RGBColor(0x8A, 0x66, 0x0E)
FG = RGBColor(0xE5, 0xE5, 0xE5)
MUTED = RGBColor(0x8A, 0x8A, 0x8A)
DIM = RGBColor(0x55, 0x55, 0x55)
BORDER = RGBColor(0x2A, 0x2A, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE0, 0x50, 0x50)
BLUE = RGBColor(0x4A, 0x90, 0xD9)
GREEN = RGBColor(0x6C, 0xB5, 0x6C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ARCHETYPE_DIR = "public/archetypes"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ============================================================
# Helpers
# ============================================================
def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def tx(slide, left, top, width, height, text, *, size=13, color=FG,
       bold=False, align=PP_ALIGN.LEFT, font="Courier New", spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_after = Pt(spacing)
    return tf


def bullets(slide, left, top, width, height, items, *, size=12, color=FG):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Courier New"
        p.space_after = Pt(5)
    return tf


def rect(slide, left, top, width, height, *, fill=CARD, border=BORDER,
         border_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = border
        s.line.width = Pt(border_w)
    s.shadow.inherit = False
    return s


def accent_line(slide, left, top, width):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    return s


def slide_header(slide, kicker, title=None):
    tx(slide, Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.4),
       kicker, size=11, color=ACCENT, bold=True)
    accent_line(slide, Inches(0.8), Inches(0.9), Inches(1.8))
    if title:
        tx(slide, Inches(0.8), Inches(1.05), Inches(11.7), Inches(0.55),
           title, size=22, color=FG, bold=True)


def arrow(slide, x, y, length=Inches(0.55), color=ACCENT):
    """Horizontal chevron arrow."""
    tx(slide, x, y, length, Inches(0.4), ">", size=22, color=color, bold=True,
       align=PP_ALIGN.CENTER)


def oval(slide, left, top, diam, *, fill=ACCENT, border=ACCENT):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diam, diam)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(0.75)
    return s


def page_num(slide, n, total):
    tx(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.3),
       f"{n:02d} / {total:02d}", size=9, color=DIM, align=PP_ALIGN.RIGHT)


def footer_brand(slide):
    tx(slide, Inches(0.8), Inches(7.05), Inches(6), Inches(0.3),
       "INTERACTIVE MEMORY MACHINE  ·  v2.4  ·  TEAM 10",
       size=9, color=DIM)


TOTAL = 13  # will update if slides change


# ============================================================
# SLIDE 1 — Title
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)

# Subtle top-left scan bar
rect(s, Inches(0.8), Inches(0.8), Inches(3.2), Pt(1.5),
     fill=ACCENT, border=None)
tx(s, Inches(0.8), Inches(0.9), Inches(6), Inches(0.3),
   "TRANSMISSION // SUBCONSCIOUS CHANNEL", size=10, color=ACCENT, bold=True)

tx(s, Inches(0.8), Inches(2.3), Inches(12), Inches(1.4),
   "INTERACTIVE MEMORY MACHINE", size=44, color=FG, bold=True)
tx(s, Inches(0.8), Inches(3.5), Inches(12), Inches(0.6),
   "v2.4 // NODE-BASED VISUAL PROGRAMMING + 2D-3D PIPELINE",
   size=16, color=ACCENT, bold=True)

accent_line(s, Inches(0.8), Inches(4.3), Inches(2.5))

tx(s, Inches(0.8), Inches(4.5), Inches(10), Inches(0.4),
   "Creative Machine Learning // Term 6 // SUTD", size=13, color=MUTED)
tx(s, Inches(0.8), Inches(4.95), Inches(10), Inches(0.4),
   "Team 10   JW  /  NT  /  VP  /  ZH", size=13, color=MUTED)

# Bottom meta strip
rect(s, Inches(0.8), Inches(6.5), Inches(11.7), Pt(1), fill=BORDER, border=None)
tx(s, Inches(0.8), Inches(6.6), Inches(6), Inches(0.3),
   "STATUS  ·  PRE-CRIT POLISH", size=10, color=MUTED)
tx(s, Inches(7), Inches(6.6), Inches(5.5), Inches(0.3),
   "BUILD 2026-04-20   CRIT 2026-04-21",
   size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 2 — The Concept
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "01 // CONCEPT", "Glitch = meaningful distortion of spatial recall")

# Left — claim
tx(s, Inches(0.8), Inches(1.9), Inches(6), Inches(0.8),
   "The glitch is not aesthetic.",
   size=26, color=FG, bold=True, font="Georgia")
tx(s, Inches(0.8), Inches(2.8), Inches(6), Inches(0.6),
   "It is how the subconscious warps architecture.",
   size=16, color=ACCENT, font="Georgia")

bullets(s, Inches(0.8), Inches(3.7), Inches(6), Inches(3),
        [
            "  Memory is never a photograph.",
            "  It omits, compresses, tilts, and fuses.",
            "",
            "  We map these distortions onto form:",
            "    -  scored parameters  ->  geometric operations",
            "    -  tone archetypes    ->  starter shape-language",
            "    -  subject's recall   ->  their own sculpture",
        ], size=13, color=MUTED)

# Right — visual: clean cube -> distorted cube
box_x = Inches(7.5); box_y = Inches(1.9); box_w = Inches(5); box_h = Inches(4.5)
rect(s, box_x, box_y, box_w, box_h, fill=CARD, border=BORDER)
tx(s, box_x + Inches(0.3), box_y + Inches(0.2), box_w - Inches(0.6), Inches(0.3),
   "BEFORE  ->  AFTER", size=10, color=ACCENT, bold=True)

# "BEFORE" — clean cube (isometric triangle approximation using rectangles)
cx = box_x + Inches(0.6); cy = box_y + Inches(1.1)
rect(s, cx, cy, Inches(1.4), Inches(1.4), fill=CARD_HI, border=MUTED)
rect(s, cx + Inches(0.2), cy - Inches(0.2), Inches(1.4), Inches(0.2),
     fill=ACCENT_DIM, border=MUTED)
tx(s, cx, cy + Inches(1.6), Inches(1.4), Inches(0.3),
   "ORDERED", size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Arrow
tx(s, box_x + Inches(2.1), cy + Inches(0.5), Inches(0.7), Inches(0.5),
   "->", size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# "AFTER" — fractured cube (scattered rectangles)
dx = box_x + Inches(2.9); dy = box_y + Inches(0.9)
rect(s, dx, dy, Inches(0.8), Inches(0.9), fill=CARD_HI, border=ACCENT_DIM)
rect(s, dx + Inches(0.9), dy + Inches(0.3), Inches(0.6), Inches(0.7), fill=CARD_HI, border=ACCENT_DIM)
rect(s, dx + Inches(0.3), dy + Inches(1.1), Inches(1.1), Inches(0.5), fill=CARD_HI, border=ACCENT)
rect(s, dx + Inches(1.3), dy + Inches(1.2), Inches(0.4), Inches(0.4), fill=ACCENT_DIM, border=ACCENT_DIM)
rect(s, dx - Inches(0.1), dy + Inches(0.3), Inches(0.2), Inches(0.2), fill=ACCENT, border=ACCENT)
tx(s, dx, cy + Inches(1.6), Inches(1.8), Inches(0.3),
   "RECALLED", size=10, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

footer_brand(s); page_num(s, 2, TOTAL)


# ============================================================
# SLIDE 3 — Evolution Timeline
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "02 // EVOLUTION", "v1  ->  v2  ->  v2.4 — three pivots, four weeks")

timeline = [
    ("v1.0", "2026-03-28", "SCAFFOLD",
        ["Memory + dream text input",
         "Claude analysis -> image -> video",
         "End-to-end pipeline working",
         "Problem: dream framing felt whimsical"]),
    ("v2.0", "2026-04-14", "NODE PIVOT",
        ["Dropped dream framing",
         "Single subconsciousness text",
         "Node-based visual programming",
         "10 operations, Ideogram V2",
         "Morphological fusion aesthetic",
         "Instructor-validated"]),
    ("v2.2", "2026-04-18", "ART DIRECTOR",
        ["Claude rewrites mechanical prompt",
         "12 nodes (spatial / exp / appearance)",
         "7 tone archetypes as starter shapes",
         "Multi-model: Recraft V4 Pro + others"]),
    ("v2.3", "2026-04-20", "3D PIPELINE",
        ["fal.ai Trellis + Hunyuan3D V3",
         "React Three Fiber viewport",
         "GLB / OBJ / STL export",
         "Teammates 3D-printed outputs"]),
    ("v2.4", "2026-04-20", "UI POLISH",
        ["Idle/subconscious glitch effects",
         "Mechanical loaders + telemetry",
         "Wire draw-in, breath glow",
         "EB Garamond serif labels",
         "Film grain + easter eggs"]),
]

col_w = Inches(2.45)
gap = Inches(0.12)
start_x = Inches(0.8)
y = Inches(1.85)

for i, (ver, date, title, notes) in enumerate(timeline):
    x = start_x + i * (col_w + gap)
    rect(s, x, y, col_w, Inches(4.8), fill=CARD, border=BORDER)
    # Version chip
    rect(s, x + Inches(0.25), y + Inches(0.25), Inches(0.9), Inches(0.35),
         fill=ACCENT_DIM, border=ACCENT)
    tx(s, x + Inches(0.25), y + Inches(0.28), Inches(0.9), Inches(0.3),
       ver, size=11, color=FG, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(1.2), y + Inches(0.3), col_w - Inches(1.3), Inches(0.3),
       date, size=9, color=MUTED)
    tx(s, x + Inches(0.25), y + Inches(0.75), col_w - Inches(0.4), Inches(0.4),
       title, size=12, color=ACCENT, bold=True)
    bullets(s, x + Inches(0.25), y + Inches(1.2), col_w - Inches(0.4),
            Inches(3.5),
            ["- " + n for n in notes], size=9, color=FG)

tx(s, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.3),
   "Each pivot addressed a critique. The glitch was never surface.",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

footer_brand(s); page_num(s, 3, TOTAL)


# ============================================================
# SLIDE 4 — Full Pipeline (10 steps as flow diagram)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "03 // PIPELINE", "10 stages — from recalled text to exported mesh")

pipeline = [
    ("01", "INPUT", "Subject ID +\nrecall text\n(<=600 chars)", FG),
    ("02", "ANALYZE", "Claude Haiku\nkeywords +\n12 scores +\narchetype", ACCENT),
    ("03", "CANVAS", "Drag chips\nWire to 12\noperation\nnodes", FG),
    ("04", "DIRECT", "Claude rewrites\nmechanical prompt\n-> visual brief", ACCENT),
    ("05", "GENERATE", "fal.ai\nRecraft V4 Pro\n(or Nano/Ideogram)", FG),
    ("06", "NAME", "Claude returns\n{name,\nexplanation}\nJSON", ACCENT),
    ("07", "VIEW 2D", "Image + title +\neditable\nexplanation", FG),
    ("08", "RECONSTRUCT", "fal.ai Trellis\nor Hunyuan3D\n-> GLB mesh", ACCENT),
    ("09", "ORBIT", "React Three Fiber\nSTUDIO / MOODY\nlighting", FG),
    ("10", "EXPORT", "GLB native +\nOBJ + STL\nclient-side", ACCENT),
]

# 2 rows of 5
row_y = [Inches(1.9), Inches(4.45)]
cell_w = Inches(2.3)
cell_h = Inches(2.1)
cell_gap = Inches(0.2)

for i, (num, title, desc, color) in enumerate(pipeline):
    row = i // 5
    col = i % 5
    x = Inches(0.8) + col * (cell_w + cell_gap)
    y = row_y[row]

    border_color = ACCENT if color == ACCENT else BORDER
    rect(s, x, y, cell_w, cell_h, fill=CARD, border=border_color)

    # Step number
    tx(s, x + Inches(0.2), y + Inches(0.15), Inches(1), Inches(0.4),
       num, size=18, color=color, bold=True)
    # Title
    tx(s, x + Inches(0.2), y + Inches(0.7), cell_w - Inches(0.4), Inches(0.35),
       title, size=12, color=FG, bold=True)
    # Desc
    tx(s, x + Inches(0.2), y + Inches(1.1), cell_w - Inches(0.4), Inches(1),
       desc, size=9, color=MUTED, spacing=1.0)

    # Arrow between cells in same row (except last in row)
    if col < 4:
        ax = x + cell_w + Inches(0.0)
        ay = y + Inches(0.75)
        tx(s, ax, ay, cell_gap, Inches(0.4),
           ">", size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Vertical turn arrow from slot 05 -> slot 06 (end of row 1 -> start of row 2)
tx(s, Inches(11.8), Inches(4.0), Inches(0.8), Inches(0.4),
   "v", size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
tx(s, Inches(0.3), Inches(4.0), Inches(0.5), Inches(0.4),
   "v", size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Legend
tx(s, Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.3),
   "WHITE = user action   ·   GOLD = model inference (Claude / fal.ai)",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

footer_brand(s); page_num(s, 4, TOTAL)


# ============================================================
# SLIDE 5 — 7 Tone Archetypes (with ref images)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "04 // ARCHETYPES", "7 tonal starters — shape-language from mood")

archetypes = [
    ("ORGANIC",     "organic.jpg",     "Calm, nostalgic",      "Smooth biomorphic curves"),
    ("CRYSTALLINE", "crystalline.jpg", "Tense, anxious",       "Angular faceted planes"),
    ("TWISTED",     "twisted.jpg",     "Disorienting, dynamic","Coiling spirals"),
    ("SKELETAL",    "skeletal.jpg",    "Vulnerable, exposed",  "Thin branching members"),
    ("MONOLITHIC",  "monolithic.jpg",  "Oppressive, heavy",    "Dense compact mass"),
    ("FRAGMENTED",  "fragmented.jpg",  "Chaotic, traumatic",   "Scattered shards"),
    ("NESTED",      "nested.jpg",      "Introspective",        "Concentric shells"),
]

# 7 cards in one row, 1.65 wide each
card_w = Inches(1.65)
card_h = Inches(4.5)
gap = Inches(0.12)
x0 = Inches(0.8)
y0 = Inches(1.85)

for i, (name, img, mood, form) in enumerate(archetypes):
    x = x0 + i * (card_w + gap)
    rect(s, x, y0, card_w, card_h, fill=CARD, border=BORDER)

    img_path = os.path.join(ARCHETYPE_DIR, img)
    if os.path.exists(img_path):
        try:
            s.shapes.add_picture(img_path,
                                 x + Inches(0.1), y0 + Inches(0.1),
                                 width=card_w - Inches(0.2),
                                 height=Inches(2.0))
        except Exception as e:
            tx(s, x + Inches(0.15), y0 + Inches(0.8), card_w - Inches(0.3), Inches(0.4),
               "[img]", size=10, color=MUTED, align=PP_ALIGN.CENTER)
    else:
        rect(s, x + Inches(0.1), y0 + Inches(0.1),
             card_w - Inches(0.2), Inches(2.0),
             fill=CARD_HI, border=BORDER)
        tx(s, x + Inches(0.15), y0 + Inches(1.0),
           card_w - Inches(0.3), Inches(0.3),
           "[ref]", size=9, color=MUTED, align=PP_ALIGN.CENTER)

    tx(s, x + Inches(0.1), y0 + Inches(2.25), card_w - Inches(0.2), Inches(0.35),
       name, size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(0.1), y0 + Inches(2.65), card_w - Inches(0.2), Inches(0.8),
       mood, size=9, color=FG, align=PP_ALIGN.CENTER)
    tx(s, x + Inches(0.1), y0 + Inches(3.35), card_w - Inches(0.2), Inches(1),
       form, size=9, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.0)

tx(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
   "Claude classifies the recall text into exactly one archetype.",
   size=11, color=FG, align=PP_ALIGN.CENTER)
tx(s, Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.3),
   "The archetype seeds the starter shape; the 12 nodes deform it.",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

page_num(s, 5, TOTAL)


# ============================================================
# SLIDE 6 — 12 Operation Nodes
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "05 // OPERATIONS", "12 nodes — three families of distortion")

# Three column groups
groups = [
    ("SPATIAL", "unipolar 0-1",
        [
            ("CLARITY",        "Boundary Diffusion",      "Edges dissolve into neighbours"),
            ("COMPLETENESS",   "Mass Subtraction",        "Voids and hollows carve volume"),
            ("STABILITY",      "Gravitational Shift",     "Tilts and suspends off-axis"),
            ("MISASSOC.",      "Cross-Morph Fusion",      "Incompatible shapes merge"),
            ("VULNERABILITY",  "Interior Reveal",         "Peels surfaces, exposes inside"),
            ("INTIMACY",       "Scale Compression",       "Mass folds inward, crushing"),
        ]),
    ("EXPERIENCE", "bipolar, 0.5 neutral",
        [
            ("TEMPERATURE",    "Freeze  <->  Melt",       "Ice/frost  vs  molten flow"),
            ("PRESSURE",       "Smooth  <->  Shatter",    "Glassy polish vs jagged spikes"),
            ("LUMINOSITY",     "Shadow  <->  Radiance",   "Consumed dark vs blinding glow"),
        ]),
    ("APPEARANCE", "unipolar 0-1",
        [
            ("MATERIAL",       "Substance Identity",      "Connected keywords set material"),
            ("TEXTURE",        "Surface Grain",           "Roughness / smoothness"),
            ("COLOR",          "Chromatic Palette",       "Palette + saturation"),
        ]),
]

col_x = [Inches(0.8), Inches(5.3), Inches(9.3)]
col_w = [Inches(4.3), Inches(3.8), Inches(3.7)]
y0 = Inches(1.85)

for ci, (group_name, group_sub, nodes) in enumerate(groups):
    x = col_x[ci]; w = col_w[ci]
    # Group header
    rect(s, x, y0, w, Inches(0.55), fill=ACCENT_DIM, border=ACCENT)
    tx(s, x + Inches(0.2), y0 + Inches(0.1), w - Inches(0.4), Inches(0.3),
       group_name, size=12, color=FG, bold=True)
    tx(s, x + Inches(0.2), y0 + Inches(0.33), w - Inches(0.4), Inches(0.25),
       group_sub, size=9, color=ACCENT)

    # Nodes
    for i, (name, op, effect) in enumerate(nodes):
        ny = y0 + Inches(0.7) + i * Inches(0.72)
        rect(s, x, ny, w, Inches(0.65), fill=CARD, border=BORDER)
        tx(s, x + Inches(0.2), ny + Inches(0.05), w - Inches(0.4), Inches(0.3),
           name, size=11, color=ACCENT, bold=True)
        tx(s, x + Inches(0.2), ny + Inches(0.3), w - Inches(0.4), Inches(0.25),
           op, size=9, color=FG)
        tx(s, x + Inches(0.2), ny + Inches(0.5), w - Inches(0.4), Inches(0.2),
           effect, size=8, color=MUTED)

tx(s, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.3),
   "Each node receives 0+ keyword wires; connection count + keyword content modulate intensity.",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

page_num(s, 6, TOTAL)


# ============================================================
# SLIDE 7 — Node Canvas UI (mockup)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "06 // CANVAS", "How the subject maps recall to form")

# Canvas frame
canvas_x = Inches(0.8); canvas_y = Inches(1.85)
canvas_w = Inches(11.7); canvas_h = Inches(4.6)
rect(s, canvas_x, canvas_y, canvas_w, canvas_h,
     fill=RGBColor(0x0D, 0x0D, 0x0D), border=BORDER)

# Column labels
tx(s, canvas_x + Inches(0.3), canvas_y + Inches(0.15),
   Inches(3), Inches(0.3),
   "KEYWORD FRAGMENTS", size=9, color=MUTED, bold=True)
tx(s, canvas_x + Inches(8.5), canvas_y + Inches(0.15),
   Inches(3), Inches(0.3),
   "OPERATION NODES", size=9, color=MUTED, bold=True)

# Keywords (left)
kws = [
    ("OBJ", "corridor"), ("OBJ", "stove"), ("PER", "mother"),
    ("SPA", "narrow"),   ("EXP", "cold"),  ("EVT", "shifting"),
    ("OBJ", "kettle"),   ("SPA", "high"),
]
cat_colors = {"OBJ": FG, "PER": ACCENT, "EVT": BLUE, "SPA": GREEN, "EXP": RED}
kw_x = canvas_x + Inches(0.3)
kw_y0 = canvas_y + Inches(0.6)
for i, (cat, word) in enumerate(kws):
    y = kw_y0 + i * Inches(0.48)
    rect(s, kw_x, y, Inches(2.4), Inches(0.4), fill=CARD, border=BORDER)
    tx(s, kw_x + Inches(0.1), y + Inches(0.08), Inches(0.4), Inches(0.25),
       cat, size=8, color=cat_colors.get(cat, FG), bold=True)
    tx(s, kw_x + Inches(0.55), y + Inches(0.08), Inches(1.7), Inches(0.25),
       word, size=10, color=FG)
    # output port
    oval(s, kw_x + Inches(2.35), y + Inches(0.13), Inches(0.15),
         fill=ACCENT, border=ACCENT)

# Operation nodes (right)
ops = [
    ("CLARITY",        0.72),
    ("COMPLETENESS",   0.35),
    ("STABILITY",      0.81),
    ("MISASSOCIATION", 0.28),
    ("VULNERABILITY",  0.55),
    ("INTIMACY",       0.67),
]
op_x = canvas_x + Inches(8.5)
op_y0 = canvas_y + Inches(0.6)
op_w = Inches(2.9)
for i, (name, score) in enumerate(ops):
    y = op_y0 + i * Inches(0.62)
    rect(s, op_x, y, op_w, Inches(0.54), fill=CARD, border=BORDER)
    # input port
    oval(s, op_x - Inches(0.08), y + Inches(0.2), Inches(0.17),
         fill=ACCENT, border=ACCENT)
    tx(s, op_x + Inches(0.15), y + Inches(0.05),
       op_w - Inches(0.8), Inches(0.25),
       name, size=10, color=ACCENT, bold=True)
    tx(s, op_x + op_w - Inches(0.65), y + Inches(0.05),
       Inches(0.6), Inches(0.25),
       f"{score:.2f}", size=10, color=FG, align=PP_ALIGN.RIGHT)
    # score bar
    bar_bg_w = Inches(2.5)
    rect(s, op_x + Inches(0.15), y + Inches(0.35), bar_bg_w, Inches(0.08),
         fill=BORDER, border=BORDER)
    rect(s, op_x + Inches(0.15), y + Inches(0.35),
         Inches(2.5 * score), Inches(0.08),
         fill=ACCENT, border=ACCENT)

# Connection wires (center) — connectors from kw port to op port
connections = [(0, 0), (1, 5), (3, 2), (4, 5), (5, 3), (6, 0), (7, 2)]
for (ki, oi) in connections:
    start_x = kw_x + Inches(2.42)
    start_y = kw_y0 + ki * Inches(0.48) + Inches(0.2)
    end_x = op_x - Inches(0.02)
    end_y = op_y0 + oi * Inches(0.62) + Inches(0.28)
    c = s.shapes.add_connector(MSO_CONNECTOR.CURVE,
                               start_x, start_y, end_x, end_y)
    c.line.color.rgb = ACCENT
    c.line.width = Pt(1.25)

# Center caption
tx(s, canvas_x + Inches(3.3), canvas_y + Inches(2.15),
   Inches(5), Inches(0.4),
   "DRAG  ->  WIRE  ->  DEFORM",
   size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
tx(s, canvas_x + Inches(3.3), canvas_y + Inches(2.55),
   Inches(5), Inches(0.6),
   "The wiring pattern itself is the\nsubject's testimony of their memory.",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Hint strip
tx(s, Inches(0.8), Inches(6.65), Inches(11.7), Inches(0.3),
   "PAN  middle-drag   ·   ZOOM  wheel   ·   DELETE  click wire   ·   MULTI-WIRE  per node allowed",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

page_num(s, 7, TOTAL)


# ============================================================
# SLIDE 8 — Art Director pipeline
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "07 // ART DIRECTOR", "Why Claude rewrites the prompt before generation")

# LEFT: mechanical prompt
mech_x = Inches(0.8); mech_y = Inches(1.85)
mech_w = Inches(4.0); mech_h = Inches(4.8)
rect(s, mech_x, mech_y, mech_w, mech_h, fill=CARD, border=BORDER)
tx(s, mech_x + Inches(0.2), mech_y + Inches(0.15), mech_w - Inches(0.4), Inches(0.3),
   "MECHANICAL PROMPT", size=10, color=ACCENT, bold=True)
tx(s, mech_x + Inches(0.2), mech_y + Inches(0.45), mech_w - Inches(0.4), Inches(0.3),
   "(from connections + scores)", size=8, color=MUTED)
bullets(s, mech_x + Inches(0.2), mech_y + Inches(0.85),
        mech_w - Inches(0.4), mech_h - Inches(1),
        [
            "archetype: CRYSTALLINE",
            "",
            "clarity 0.72",
            "  -> boundaries dissolve",
            "completeness 0.35",
            "  -> heavy mass subtraction",
            "stability 0.81",
            "  -> mild tilt",
            "",
            "kw(corridor, narrow, high)",
            "  -> spatial nodes",
            "kw(cold, shifting)",
            "  -> temperature low",
            "",
            "...+ 9 more parameters",
        ], size=10, color=FG)

# ARROW 1
tx(s, mech_x + mech_w + Inches(0.05), Inches(3.8), Inches(0.55), Inches(0.5),
   "->", size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# MIDDLE: Claude brain
cl_x = Inches(5.45); cl_y = Inches(1.85)
cl_w = Inches(2.8); cl_h = Inches(4.8)
rect(s, cl_x, cl_y, cl_w, cl_h, fill=CARD_HI, border=ACCENT, border_w=1.5)
tx(s, cl_x + Inches(0.2), cl_y + Inches(0.15), cl_w - Inches(0.4), Inches(0.3),
   "ART DIRECTOR", size=10, color=ACCENT, bold=True)
tx(s, cl_x + Inches(0.2), cl_y + Inches(0.45), cl_w - Inches(0.4), Inches(0.3),
   "Claude Haiku", size=10, color=FG, bold=True)
tx(s, cl_x + Inches(0.2), cl_y + Inches(0.75), cl_w - Inches(0.4), Inches(0.3),
   "context-aware", size=8, color=MUTED)

bullets(s, cl_x + Inches(0.2), cl_y + Inches(1.2),
        cl_w - Inches(0.4), cl_h - Inches(1.4),
        [
            "Reads:",
            "  + original recall",
            "  + archetype",
            "  + all 12 scores",
            "  + wiring pattern",
            "",
            "Writes:",
            "  single cohesive",
            "  visual brief for",
            "  text-to-image model",
            "",
            "Physically plausible",
            "sculpture, clean",
            "background.",
        ], size=9, color=FG)

# ARROW 2
tx(s, cl_x + cl_w + Inches(0.05), Inches(3.8), Inches(0.55), Inches(0.5),
   "->", size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# RIGHT: visual brief
vb_x = Inches(8.9); vb_y = Inches(1.85)
vb_w = Inches(3.6); vb_h = Inches(4.8)
rect(s, vb_x, vb_y, vb_w, vb_h, fill=CARD, border=BORDER)
tx(s, vb_x + Inches(0.2), vb_y + Inches(0.15), vb_w - Inches(0.4), Inches(0.3),
   "VISUAL BRIEF", size=10, color=ACCENT, bold=True)
tx(s, vb_x + Inches(0.2), vb_y + Inches(0.45), vb_w - Inches(0.4), Inches(0.3),
   "(sent to fal.ai)", size=8, color=MUTED)

tx(s, vb_x + Inches(0.2), vb_y + Inches(0.85), vb_w - Inches(0.4), vb_h - Inches(1.0),
   "A crystalline sculpture of\n"
   "a narrow corridor, faceted\n"
   "ice-like planes dissolving\n"
   "at the edges, carved with\n"
   "deep hollows where mass\n"
   "has been lost. Tilts gently\n"
   "off-axis. Cold palette, matte\n"
   "frost texture. Studio lighting,\n"
   "clean white background,\n"
   "physically plausible, single\n"
   "sculptural object.",
   size=10, color=FG, font="Georgia", spacing=1.5)

tx(s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.3),
   "Without the art director, text-to-image models treat mechanical specs as abstract noise.",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

page_num(s, 8, TOTAL)


# ============================================================
# SLIDE 9 — 2D -> 3D Reconstruction
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "08 // 2D  ->  3D", "From generated image to holdable mesh")

# Left — reconstruction flow (4 stages)
stages = [
    ("IMAGE",      "Generated 2D\nartifact", FG),
    ("RECONSTRUCT","Trellis or\nHunyuan3D V3\nvia fal.ai",  ACCENT),
    ("MESH",       "GLB with\ngeometry +\nmaterials",   FG),
    ("VIEWER",     "React Three\nFiber viewport",      ACCENT),
]
sx0 = Inches(0.8); sy = Inches(1.85)
sw = Inches(2.6); sh = Inches(2.2); sgap = Inches(0.15)

for i, (title, desc, color) in enumerate(stages):
    x = sx0 + i * (sw + sgap)
    border_c = ACCENT if color == ACCENT else BORDER
    rect(s, x, sy, sw, sh, fill=CARD, border=border_c)
    tx(s, x + Inches(0.2), sy + Inches(0.25), sw - Inches(0.4), Inches(0.4),
       title, size=13, color=color, bold=True)
    tx(s, x + Inches(0.2), sy + Inches(0.85), sw - Inches(0.4), Inches(1.2),
       desc, size=10, color=MUTED, spacing=1.2)
    if i < 3:
        tx(s, x + sw - Inches(0.05), sy + Inches(0.8), sgap + Inches(0.1), Inches(0.4),
           ">", size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Lower section — model comparison table
ty = Inches(4.35)
rect(s, Inches(0.8), ty, Inches(11.7), Inches(2.2), fill=CARD, border=BORDER)
tx(s, Inches(1.0), ty + Inches(0.15), Inches(11.3), Inches(0.3),
   "MODEL TOGGLE — chosen per-generation via UI", size=10, color=ACCENT, bold=True)

# Header row
hdr_y = ty + Inches(0.6)
headers = [("MODEL", Inches(1.0), Inches(2.2)),
           ("fal ID", Inches(3.2), Inches(3.3)),
           ("COST / GEN", Inches(6.5), Inches(1.5)),
           ("FIDELITY", Inches(8.0), Inches(1.7)),
           ("WHEN TO USE", Inches(9.7), Inches(2.8))]
for (h, hx, hw) in headers:
    tx(s, hx, hdr_y, hw, Inches(0.3), h, size=10, color=MUTED, bold=True)

rect(s, Inches(1.0), hdr_y + Inches(0.35), Inches(11.3), Pt(0.75),
     fill=BORDER, border=None)

# Row 1 — Trellis
r1y = hdr_y + Inches(0.55)
tx(s, Inches(1.0), r1y, Inches(2.2), Inches(0.3),
   "TRELLIS", size=11, color=ACCENT, bold=True)
tx(s, Inches(3.2), r1y, Inches(3.3), Inches(0.3),
   "fal-ai/trellis", size=10, color=FG)
tx(s, Inches(6.5), r1y, Inches(1.5), Inches(0.3),
   "~$0.02", size=10, color=GREEN)
tx(s, Inches(8.0), r1y, Inches(1.7), Inches(0.3),
   "standard", size=10, color=FG)
tx(s, Inches(9.7), r1y, Inches(2.8), Inches(0.3),
   "default, fast iteration", size=10, color=MUTED)

# Row 2 — Hunyuan3D
r2y = r1y + Inches(0.55)
tx(s, Inches(1.0), r2y, Inches(2.2), Inches(0.3),
   "HUNYUAN3D V3", size=11, color=ACCENT, bold=True)
tx(s, Inches(3.2), r2y, Inches(3.3), Inches(0.3),
   "fal-ai/hunyuan3d-v3/image-to-3d", size=10, color=FG)
tx(s, Inches(6.5), r2y, Inches(1.5), Inches(0.3),
   "~$0.16", size=10, color=RED)
tx(s, Inches(8.0), r2y, Inches(1.7), Inches(0.3),
   "PBR materials", size=10, color=FG)
tx(s, Inches(9.7), r2y, Inches(2.8), Inches(0.3),
   "crit / exhibition hero",
   size=10, color=MUTED)

tx(s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.3),
   "Exports: GLB (native from fal) + OBJ + STL (client-side via three-stdlib for 3D-print slicers).",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

page_num(s, 9, TOTAL)


# ============================================================
# SLIDE 10 — UI Theme Duality (mechanical / subconscious)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "09 // UI THEME", "Mechanical precision  <->  subconscious decay")

# Two vertical halves
# LEFT: ACTIVE
lx = Inches(0.8); ly = Inches(1.85); lw = Inches(5.75); lh = Inches(4.85)
rect(s, lx, ly, lw, lh, fill=CARD, border=ACCENT)
tx(s, lx + Inches(0.3), ly + Inches(0.25), lw - Inches(0.6), Inches(0.4),
   "ACTIVE / MECHANICAL", size=14, color=ACCENT, bold=True)
tx(s, lx + Inches(0.3), ly + Inches(0.7), lw - Inches(0.6), Inches(0.3),
   "input + mapping + generating", size=10, color=MUTED)

bullets(s, lx + Inches(0.3), ly + Inches(1.2), lw - Inches(0.6), Inches(3.6),
        [
            "+ Scan bars / elapsed counters",
            "+ Telemetry glyphs + frame counter",
            "+ Wire draw-in animation",
            "+ Node breath glow (connected)",
            "+ Badge pop on connect",
            "+ EB Garamond serif for artifact name",
            "+ Courier Mono for data / code",
            "+ Precise, deliberate, robotic",
        ], size=12, color=FG)

# RIGHT: IDLE
rx = Inches(6.75); ry = ly; rw = Inches(5.75); rh = lh
rect(s, rx, ry, rw, rh, fill=CARD, border=MUTED)
tx(s, rx + Inches(0.3), ry + Inches(0.25), rw - Inches(0.6), Inches(0.4),
   "IDLE / SUBCONSCIOUS", size=14, color=FG, bold=True)
tx(s, rx + Inches(0.3), ry + Inches(0.7), rw - Inches(0.6), Inches(0.3),
   "input screen only, >=12s no input", size=10, color=MUTED)

bullets(s, rx + Inches(0.3), ry + Inches(1.2), rw - Inches(0.6), Inches(3.6),
        [
            "~ Random letter-fall on headings",
            "~ Gold scan line sweep",
            "~ Periodic TV glitch",
            "  (chromatic aberration + jitter)",
            "~ Deep glitch burst at 35s (once)",
            "~ Film grain overlay always on (~6%)",
            "~ Halts immediately on any input",
            "~ Gated to session.status = input",
        ], size=12, color=FG)

# Bottom stripe — easter eggs
ey = Inches(6.85)
tx(s, Inches(0.8), ey, Inches(11.7), Inches(0.3),
   "EASTER EGGS  ·  Konami = DEV preset  ·  Midnight (00-04h) = NIGHT SHIFT subtitle  ·  35s idle = deep burst",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

page_num(s, 10, TOTAL)


# ============================================================
# SLIDE 11 — Tech Stack
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "10 // STACK", "Everything the machine runs on")

stack_rows = [
    ("FRAMEWORK",  "Next.js 15 (App Router) + TypeScript + Tailwind CSS"),
    ("UI",         "React 19  ·  useReducer  ·  pointer events + SVG bezier wires"),
    ("3D VIEWER",  "React Three Fiber + drei (OrbitControls, Bounds, Grid, GizmoHelper, Environment)"),
    ("EXPORTS",    "GLB (native)  ·  OBJ + STL (three-stdlib client-side)"),
    ("CLAUDE API", "Haiku (dev / art director / naming)  ·  Sonnet (analysis for production)"),
    ("TEXT->IMG",  "fal.ai  ·  Recraft V4 Pro (default)  ·  Nano Banana Pro  ·  Ideogram V2"),
    ("IMG->3D",    "fal.ai  ·  Trellis (default, ~$0.02)  ·  Hunyuan3D V3 PBR (~$0.16)"),
    ("HOSTING",    "Vercel (pre-deploy)  ·  .env.local for keys (gitignored)"),
    ("REPO",       "github.com/ZawYeHtet2001/Glitch_CML_project  ·  branch: master"),
]

y = Inches(1.85)
for i, (label, val) in enumerate(stack_rows):
    rect(s, Inches(0.8), y, Inches(11.7), Inches(0.48),
         fill=CARD if i % 2 == 0 else CARD_HI, border=BORDER)
    tx(s, Inches(1.0), y + Inches(0.11), Inches(2.5), Inches(0.3),
       label, size=11, color=ACCENT, bold=True)
    tx(s, Inches(3.6), y + Inches(0.11), Inches(8.8), Inches(0.3),
       val, size=11, color=FG)
    y += Inches(0.53)

tx(s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.3),
   "Rule of thumb: Claude for meaning, fal.ai for pixels / meshes, local JS for export.",
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

page_num(s, 11, TOTAL)


# ============================================================
# SLIDE 12 — Status + Final Crit Plan
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "11 // STATUS", "Ready state + last-mile plan")

# Two columns
# LEFT: Shipped
sx = Inches(0.8); sy = Inches(1.85); sw = Inches(5.75); sh = Inches(4.85)
rect(s, sx, sy, sw, sh, fill=CARD, border=BORDER)
tx(s, sx + Inches(0.3), sy + Inches(0.2), sw - Inches(0.4), Inches(0.4),
   "SHIPPED", size=13, color=GREEN, bold=True)

bullets(s, sx + Inches(0.3), sy + Inches(0.75), sw - Inches(0.4), sh - Inches(1),
        [
            "[x] v1 pipeline (2026-03-28)",
            "[x] Node-based UI pivot (2026-04-14)",
            "[x] Art director + 12 nodes (04-18)",
            "[x] 7 tone archetypes (04-19)",
            "[x] Multi-model 2D (Recraft/Nano/Ideogram)",
            "[x] Trellis + Hunyuan3D 2D->3D (04-20)",
            "[x] R3F viewer + GLB/OBJ/STL export",
            "[x] Editable artifact names",
            "[x] Mechanical loaders + telemetry",
            "[x] Idle / TV glitch effects",
            "[x] Wire draw-in + breath glow",
            "[x] Film grain + easter eggs",
            "[x] Teammates already 3D-printed",
        ], size=11, color=FG)

# RIGHT: Last-mile
rx = Inches(6.75); ry = sy; rw = sw; rh = sh
rect(s, rx, ry, rw, rh, fill=CARD, border=ACCENT)
tx(s, rx + Inches(0.3), ry + Inches(0.2), rw - Inches(0.4), Inches(0.4),
   "LAST-MILE  (MON NIGHT)", size=13, color=ACCENT, bold=True)

bullets(s, rx + Inches(0.3), ry + Inches(0.75), rw - Inches(0.4), rh - Inches(1),
        [
            "[ ] UI optimization pass (tonight)",
            "[ ] Full dry-run demo with all 4 recalls",
            "[ ] Check Claude Haiku -> Sonnet toggle",
            "[ ] Verify all 3 image models",
            "[ ] Re-export hero pieces (Hunyuan3D)",
            "[ ] Slicer test: STL opens in Bambu / Cura",
            "[ ] Deploy to Vercel (staging)",
            "[ ] Crit rehearsal — 5 min pitch",
            "",
            "FREEZE:  Tuesday morning.",
            "No new features after that.",
            "Only bug fixes + copy tweaks.",
        ], size=11, color=FG)

# Timeline strip
tx(s, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.3),
   "CRIT  ·  Tuesday 2026-04-21, evening",
   size=11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

page_num(s, 12, TOTAL)


# ============================================================
# SLIDE 13 — Team roles + Close
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
slide_header(s, "12 // TEAM", "Who owns what through crit")

members = [
    ("ZH", "Owner / Primary Dev",
        ["Architecture, pipeline, 3D",
         "UI polish + demo rehearsal",
         "This deck"]),
    ("JW", "Recall contributor",
        ["Subconsciousness text",
         "Crit support"]),
    ("NT", "Recall contributor",
        ["Subconsciousness text",
         "Crit support"]),
    ("VP", "Recall contributor",
        ["Subconsciousness text",
         "Crit support"]),
]

cx = Inches(0.8); cy = Inches(1.85)
cw = Inches(2.88); ch = Inches(3.8); cgap = Inches(0.12)

for i, (name, role, tasks) in enumerate(members):
    x = cx + i * (cw + cgap)
    border_c = ACCENT if name == "ZH" else BORDER
    rect(s, x, cy, cw, ch, fill=CARD, border=border_c)
    tx(s, x + Inches(0.3), cy + Inches(0.3), cw - Inches(0.6), Inches(0.6),
       name, size=34, color=ACCENT, bold=True)
    tx(s, x + Inches(0.3), cy + Inches(1.1), cw - Inches(0.6), Inches(0.35),
       role, size=11, color=FG, bold=True)
    bullets(s, x + Inches(0.3), cy + Inches(1.6), cw - Inches(0.6), Inches(2),
            ["- " + t for t in tasks], size=10, color=MUTED)

# Bottom CTA
cta_y = Inches(5.95)
rect(s, Inches(0.8), cta_y, Inches(11.7), Inches(1.1),
     fill=CARD_HI, border=ACCENT)
tx(s, Inches(1.0), cta_y + Inches(0.2), Inches(11.3), Inches(0.4),
   "TRY IT  ·  localhost:3004   //   REPO  ·  github.com/ZawYeHtet2001/Glitch_CML_project",
   size=12, color=ACCENT, bold=True)
tx(s, Inches(1.0), cta_y + Inches(0.6), Inches(11.3), Inches(0.4),
   "Ping ZH for Claude / fal.ai key access. All feedback welcome before Monday night freeze.",
   size=11, color=MUTED)

tx(s, Inches(0.8), Inches(7.15), Inches(11.7), Inches(0.3),
   "END  //  TRANSMISSION COMPLETE",
   size=9, color=DIM, align=PP_ALIGN.CENTER)

page_num(s, 13, TOTAL)


# ============================================================
# Save
# ============================================================
output_path = "IMM_v2_4_Team_Update.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"  {len(prs.slides)} slides, dark clinical theme (v2.4)")
