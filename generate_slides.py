"""
Generate project update PowerPoint for Interactive Memory Machine v2.
Run: python generate_slides.py
Output: IMM_v2_Update.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Colors ---
BG = RGBColor(0x0A, 0x0A, 0x0A)
CARD = RGBColor(0x14, 0x14, 0x14)
ACCENT = RGBColor(0xD4, 0xA0, 0x17)
FG = RGBColor(0xE5, 0xE5, 0xE5)
MUTED = RGBColor(0x6A, 0x6A, 0x6A)
BORDER = RGBColor(0x26, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE0, 0x50, 0x50)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=FG, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Courier New", spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(spacing)
    return tf


def add_bullet_frame(slide, left, top, width, height, items,
                     font_size=13, color=FG, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Courier New"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_rect(slide, left, top, width, height, fill_color=CARD, border_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG)

add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(1),
             "INTERACTIVE MEMORY MACHINE", font_size=36, color=ACCENT,
             bold=True, alignment=PP_ALIGN.LEFT)
add_accent_line(slide, Inches(1), Inches(3.1), Inches(3))
add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.6),
             "v2 — Node-Based Visual Programming Interface",
             font_size=18, color=FG)
add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             "Creative Machine Learning — Term 6, SUTD",
             font_size=14, color=MUTED)
add_text_box(slide, Inches(1), Inches(4.7), Inches(11), Inches(0.5),
             "Team 10: JW  |  NT  |  VP  |  ZH",
             font_size=14, color=MUTED)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "2026-04-14", font_size=12, color=MUTED)


# ============================================================
# SLIDE 2: Concept — What is Glitch?
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "THE CONCEPT", font_size=12, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(1.2),
             "Glitch is not aesthetic decoration.",
             font_size=28, color=FG, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(2),
             "It represents meaningful distortion of spatial recall — "
             "how the subconscious warps, omits, compresses, and destabilizes "
             "our memory of architectural space.",
             font_size=15, color=MUTED)

# Right side — key shift
add_rect(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(4.5))
add_text_box(slide, Inches(7.4), Inches(1.6), Inches(4.7), Inches(0.5),
             "KEY PIVOT (v1 → v2)", font_size=11, color=ACCENT, bold=True)
add_bullet_frame(slide, Inches(7.4), Inches(2.2), Inches(4.7), Inches(3.5), [
    "× Dropped 'dream' framing",
    "  Dream imagery → whimsical, not glitchy",
    "",
    "+ Single subconsciousness input",
    "  Raw spatial recall, not curated narrative",
    "",
    "+ Node-based interactive control",
    "  Subject decides how memories distort",
    "",
    "+ 6 scored parameters on a cube",
    "  Geometric operations, not aesthetics",
], font_size=12, color=FG)


# ============================================================
# SLIDE 3: The 6 Parameters
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "6 DISTORTION PARAMETERS", font_size=12, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

params = [
    ("CLARITY", "Edge Resolution / Fidelity", "Low → boundaries dissolve into fog, forms become uncertain silhouettes"),
    ("COMPLETENESS", "Subtraction / Missing Mass", "Low → carved-out voids, hollow gaps, geometry erased"),
    ("STABILITY", "Tilt / Center of Gravity", "Low → displaced from vertical, structural equilibrium broken"),
    ("MISASSOCIATION", "Collision / Hybridization", "High → objects merged wrongly, impossible material fusions"),
    ("VULNERABILITY", "Porosity / Shell Thickness", "High → surfaces perforated, interior exposed, skin peeling"),
    ("INTIMACY", "Compression / Cavity Size", "High → walls closing in, spatial scale crushed, claustrophobic"),
]

for i, (name, desc, effect) in enumerate(params):
    row = i // 2
    col = i % 2
    x = Inches(0.8) + col * Inches(6.2)
    y = Inches(1.4) + row * Inches(1.9)

    add_rect(slide, x, y, Inches(5.8), Inches(1.6))
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.2), Inches(0.4),
                 name, font_size=14, color=ACCENT, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(5.2), Inches(0.35),
                 desc, font_size=11, color=FG)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(5.2), Inches(0.5),
                 effect, font_size=10, color=MUTED)


# ============================================================
# SLIDE 4: Pipeline Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "PIPELINE ARCHITECTURE", font_size=12, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

steps = [
    ("01", "INPUT", "Subject provides ID +\nsubconsciousness text"),
    ("02", "ANALYSIS", "Claude extracts keywords\n+ scores 6 parameters"),
    ("03", "NODE MAPPING", "User wires keyword\nfragments to operations"),
    ("04", "GENERATION", "Prompt builder → fal.ai\nFlux image generation"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(1.8)

    add_rect(slide, x, y, Inches(2.7), Inches(2.5))
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(2.1), Inches(0.4),
                 num, font_size=24, color=ACCENT, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.75), Inches(2.1), Inches(0.4),
                 title, font_size=13, color=FG, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.3), Inches(2.1), Inches(1),
                 desc, font_size=11, color=MUTED)

    # Arrow between steps
    if i < 3:
        arrow_x = x + Inches(2.8)
        add_text_box(slide, arrow_x, y + Inches(1), Inches(0.3), Inches(0.4),
                     "→", font_size=18, color=ACCENT)

# Data flow detail
add_text_box(slide, Inches(0.8), Inches(5), Inches(11.5), Inches(0.4),
             "DATA FLOW", font_size=11, color=ACCENT, bold=True)
add_text_box(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.2),
             "[Text Input] → Claude API → [8-15 Keywords + 6 Scores] → "
             "[Interactive Canvas] → [User Connections] → [Prompt Builder] → "
             "[fal.ai Flux Schnell] → [Generated Image]",
             font_size=12, color=MUTED)


# ============================================================
# SLIDE 5: Node Canvas UI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "NODE CANVAS — VISUAL PROGRAMMING INTERFACE", font_size=12, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Canvas mockup background
add_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2),
         fill_color=RGBColor(0x0D, 0x0D, 0x0D), border_color=BORDER)

# Keywords (left side)
kw_label = add_text_box(slide, Inches(1.2), Inches(1.7), Inches(2), Inches(0.3),
                        "KEYWORD FRAGMENTS", font_size=9, color=MUTED)

keywords_demo = [
    ("OBJ", "corridor"), ("OBJ", "window"), ("OBJ", "stove"),
    ("PER", "mother"), ("EVT", "shifting"),
    ("SPA", "narrow"), ("SPA", "impossibly high"),
    ("EXP", "fear"), ("EXP", "cold"),
]

kw_colors = {
    "OBJ": FG, "PER": ACCENT, "EVT": RGBColor(0x8A, 0x8A, 0x8A),
    "SPA": MUTED, "EXP": RGBColor(0xB8, 0x86, 0x0B),
}

for i, (cat, text) in enumerate(keywords_demo):
    y = Inches(2.1) + i * Inches(0.42)
    add_rect(slide, Inches(1.2), y, Inches(2), Inches(0.35), fill_color=CARD)
    add_text_box(slide, Inches(1.3), y + Inches(0.03), Inches(0.5), Inches(0.3),
                 cat, font_size=8, color=kw_colors.get(cat, FG))
    add_text_box(slide, Inches(1.8), y + Inches(0.03), Inches(1.3), Inches(0.3),
                 text, font_size=10, color=FG)
    # Output port
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.15), y + Inches(0.1),
                                    Inches(0.15), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BORDER
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)

# Operations (right side)
op_label = add_text_box(slide, Inches(9), Inches(1.7), Inches(3), Inches(0.3),
                        "OPERATION NODES", font_size=9, color=MUTED)

ops_demo = [
    ("CLARITY", "0.72"), ("COMPLETENESS", "0.35"),
    ("STABILITY", "0.81"), ("MISASSOCIATION", "0.28"),
    ("VULNERABILITY", "0.55"), ("INTIMACY", "0.67"),
]

for i, (name, score) in enumerate(ops_demo):
    y = Inches(2.1) + i * Inches(0.72)
    add_rect(slide, Inches(9), y, Inches(3.2), Inches(0.6), fill_color=CARD)
    # Input port
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.85), y + Inches(0.2),
                                    Inches(0.17), Inches(0.17))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)
    add_text_box(slide, Inches(9.2), y + Inches(0.05), Inches(2), Inches(0.3),
                 name, font_size=10, color=ACCENT, bold=True)
    add_text_box(slide, Inches(11.3), y + Inches(0.05), Inches(0.7), Inches(0.3),
                 score, font_size=10, color=FG)
    # Score bar background
    bar_x = Inches(9.2)
    bar_y = y + Inches(0.4)
    add_rect(slide, bar_x, bar_y, Inches(2.2), Inches(0.08),
             fill_color=BORDER, border_color=BORDER)
    # Score bar fill
    fill_w = int(2.2 * float(score))
    add_rect(slide, bar_x, bar_y, Inches(2.2 * float(score)), Inches(0.08),
             fill_color=ACCENT, border_color=ACCENT)

# Connection lines label
add_text_box(slide, Inches(4), Inches(3.5), Inches(4), Inches(0.5),
             "← BEZIER WIRE CONNECTIONS →", font_size=11, color=ACCENT,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4), Inches(4.0), Inches(4), Inches(0.5),
             "User drags from keyword ports\nto operation input ports",
             font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)

# Interaction hints
add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4),
             "DRAG keywords freely  |  WIRE to operations  |  DELETE to remove  |  "
             "Multiple connections allowed",
             font_size=10, color=MUTED, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6: Tech Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "TECH STACK", font_size=12, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

stack_left = [
    ("FRAMEWORK", "Next.js 15 + TypeScript + Tailwind CSS"),
    ("UI ENGINE", "React 19 — pointer events + SVG wires\nNo external node-graph libraries"),
    ("STATE", "useReducer — session state + canvas state\nisolated for performance"),
    ("STYLING", "Dark clinical aesthetic\n#0A0A0A / #D4A017 / Courier New"),
]

stack_right = [
    ("ANALYSIS", "Claude API (Sonnet / Haiku)\nKeyword extraction + 6-parameter scoring"),
    ("GENERATION", "fal.ai — Flux Schnell (dev)\nFlux Pro for production"),
    ("HOSTING", "Vercel (to be deployed)"),
    ("REPOSITORY", "GitHub — ZawYeHtet2001/Glitch_CML_project"),
]

for i, (label, desc) in enumerate(stack_left):
    y = Inches(1.5) + i * Inches(1.4)
    add_rect(slide, Inches(0.8), y, Inches(5.8), Inches(1.2))
    add_text_box(slide, Inches(1.1), y + Inches(0.15), Inches(5.2), Inches(0.3),
                 label, font_size=10, color=ACCENT, bold=True)
    add_text_box(slide, Inches(1.1), y + Inches(0.5), Inches(5.2), Inches(0.7),
                 desc, font_size=11, color=FG)

for i, (label, desc) in enumerate(stack_right):
    y = Inches(1.5) + i * Inches(1.4)
    add_rect(slide, Inches(7), y, Inches(5.8), Inches(1.2))
    add_text_box(slide, Inches(7.3), y + Inches(0.15), Inches(5.2), Inches(0.3),
                 label, font_size=10, color=ACCENT, bold=True)
    add_text_box(slide, Inches(7.3), y + Inches(0.5), Inches(5.2), Inches(0.7),
                 desc, font_size=11, color=FG)


# ============================================================
# SLIDE 7: Current Status + Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "STATUS + NEXT STEPS", font_size=12, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

# Done column
add_rect(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5))
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5.2), Inches(0.4),
             "COMPLETED", font_size=11, color=ACCENT, bold=True)

done_items = [
    "[+] v1 pipeline (2026-03-28)",
    "    Full end-to-end: input → analysis → image → video",
    "",
    "[+] v2 pivot (2026-04-14)",
    "    Node-based visual programming UI",
    "    6-parameter scoring system",
    "    Keyword extraction from text",
    "    Interactive canvas with drag + wire",
    "    Prompt builder (connections → image)",
    "    Clean TypeScript build, zero errors",
    "    Pushed to GitHub",
]
add_bullet_frame(slide, Inches(1.1), Inches(2.2), Inches(5.2), Inches(4),
                 done_items, font_size=11, color=FG)

# Next column
add_rect(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5))
add_text_box(slide, Inches(7.3), Inches(1.7), Inches(5.2), Inches(0.4),
             "NEXT", font_size=11, color=ACCENT, bold=True)

next_items = [
    "[ ] Test full pipeline end-to-end",
    "    (Claude API currently overloaded)",
    "",
    "[ ] Collect narratives from all 4 members",
    "    JW, NT, VP, ZH",
    "",
    "[ ] Tune Claude prompt for better keywords",
    "",
    "[ ] Polish canvas UI",
    "    Hover effects, animations, layout",
    "",
    "[ ] 3D cube output",
    "    Parameters → geometric deformation",
    "    Physical artifact for exhibition",
    "",
    "[ ] Deploy to Vercel",
]
add_bullet_frame(slide, Inches(7.3), Inches(2.2), Inches(5.2), Inches(4),
                 next_items, font_size=11, color=FG)


# ============================================================
# SLIDE 8: Team Action Items
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
             "TEAM ACTION ITEMS", font_size=12, color=ACCENT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.0), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
             "Each team member: write a 100-200 word subconsciousness text",
             font_size=20, color=FG, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(1),
             "Describe a remembered spatial experience — a place from memory, dream, "
             "or somewhere in between. Focus on architecture, objects, people, sensations, "
             "and how the space felt. The more specific and concrete, the better the "
             "keyword extraction.",
             font_size=13, color=MUTED)

members = [
    ("JW", "Write subconsciousness text"),
    ("NT", "Write subconsciousness text"),
    ("VP", "Write subconsciousness text"),
    ("ZH", "Write text + continue dev (prompts, 3D)"),
]

for i, (name, task) in enumerate(members):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(3.8)
    add_rect(slide, x, y, Inches(2.7), Inches(1.5))
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(2.1), Inches(0.5),
                 name, font_size=22, color=ACCENT, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(2.1), Inches(0.6),
                 task, font_size=11, color=FG)

add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
             "Try it: http://localhost:3000  |  Repo: github.com/ZawYeHtet2001/Glitch_CML_project",
             font_size=12, color=MUTED, alignment=PP_ALIGN.CENTER)


# ============================================================
# Save
# ============================================================
output_path = "IMM_v2_Update.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"  8 slides, dark clinical theme")
